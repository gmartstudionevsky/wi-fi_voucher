from __future__ import annotations

import shutil
from pathlib import Path
import copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from pypdf import PdfMerger

PASSWORD_TOKEN = "{{PASSWORD}}"
QR_TOKEN = "{{QR_WIFI}}"

def _iter_shapes_recursive(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(sh.shapes)

def _find_textbox(slide, token: str):
    for sh in _iter_shapes_recursive(slide.shapes):
        if sh.has_text_frame and token in sh.text_frame.text:
            return sh
    return None

def _remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)

def _replace_text_in_runs(text_frame, token: str, replacement: str) -> bool:
    replaced = False
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if token in run.text:
                run.text = run.text.replace(token, replacement)
                replaced = True
        if replaced:
            continue
        full_text = "".join(run.text for run in paragraph.runs)
        if token not in full_text:
            continue
        if paragraph.runs:
            first_run = paragraph.runs[0]
            for run in list(paragraph.runs)[1:]:
                run._element.getparent().remove(run._element)
            first_run.text = full_text.replace(token, replacement)
            replaced = True
    return replaced

def _replace_password(slide, password: str):
    for sh in _iter_shapes_recursive(slide.shapes):
        if not sh.has_text_frame:
            continue
        if PASSWORD_TOKEN in sh.text_frame.text:
            _replace_text_in_runs(sh.text_frame, PASSWORD_TOKEN, password)

def _insert_qr(slide, qr_png_path: str):
    box = _find_textbox(slide, QR_TOKEN)
    if box:
        left, top, width, height = box.left, box.top, box.width, box.height
        _remove_shape(box)
        slide.shapes.add_picture(qr_png_path, left, top, width=width, height=height)
        return True

    # fallback (если маркер удалили): ставим под “Пароль/Password”
    label = None
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip().lower()
            if t.startswith("пароль") or t.startswith("password"):
                label = sh
                break
    if not label:
        return False

    left = label.left
    top = label.top + label.height + Emu(120000)
    size = Emu(1150000)
    slide.shapes.add_picture(qr_png_path, left, top, width=size, height=size)
    return True

def render_single_brochure_pptx(template_path: str, password: str, qr_png_path: str, out_pptx_path: str):
    # Копируем шаблон как есть — вся графика/лого/QR внутри сохраняются
    shutil.copy(template_path, out_pptx_path)

    prs = Presentation(out_pptx_path)
    for slide in prs.slides:
        _replace_password(slide, password)
        _insert_qr(slide, qr_png_path)
    prs.save(out_pptx_path)

def _copy_slide(dest_prs: Presentation, src_slide):
    layout = dest_prs.slide_layouts[0]
    new_slide = dest_prs.slides.add_slide(layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    for shape in src_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(copy.deepcopy(shape._element), "p:extLst")
    for rel in src_slide.part.rels:
        if "notesSlide" in rel.reltype or rel.reltype.endswith("/slideLayout"):
            continue
        new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
    return new_slide

def build_pptx(
    template_ru: str,
    template_en: str,
    ru_passwords: list[str],
    en_passwords: list[str],
    qr_png_paths: list[str],
    out_pptx_path: str
):
    ru_template = Presentation(template_ru)
    en_template = Presentation(template_en)

    out_prs = Presentation()
    out_prs.slide_width = ru_template.slide_width
    out_prs.slide_height = ru_template.slide_height

    qr_idx = 0

    for pwd in ru_passwords:
        qr_path = qr_png_paths[qr_idx]
        qr_idx += 1
        for src_slide in ru_template.slides:
            slide = _copy_slide(out_prs, src_slide)
            _replace_password(slide, pwd)
            _insert_qr(slide, qr_path)

    for pwd in en_passwords:
        qr_path = qr_png_paths[qr_idx]
        qr_idx += 1
        for src_slide in en_template.slides:
            slide = _copy_slide(out_prs, src_slide)
            _replace_password(slide, pwd)
            _insert_qr(slide, qr_path)

    out_prs.save(out_pptx_path)

def convert_pptx_to_pdf(soffice_bin: str, pptx_path: str, out_dir: str) -> str:
    import subprocess
    pptx_path = str(Path(pptx_path).resolve())
    out_dir = str(Path(out_dir).resolve())
    cmd = [
        soffice_bin,
        "--headless", "--nologo", "--nofirststartwizard", "--norestore",
        "--convert-to", "pdf",
        "--outdir", out_dir,
        pptx_path
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f"LibreOffice convert failed ({proc.returncode}): {proc.stderr.strip() or proc.stdout.strip()}")
    pdf_path = str(Path(out_dir) / (Path(pptx_path).stem + ".pdf"))
    if not Path(pdf_path).exists():
        raise RuntimeError("PDF file was not produced by LibreOffice.")
    return pdf_path

def build_merged_pdf(
    soffice_bin: str,
    template_ru: str,
    template_en: str,
    ru_passwords: list[str],
    en_passwords: list[str],
    qr_png_paths: list[str],
    work_dir: str,
    out_pdf_path: str
):
    work = Path(work_dir)
    pdf_dir = work / "pdf_parts"
    pdf_dir.mkdir(parents=True, exist_ok=True)

    pdfs: list[str] = []

    # QR paths aligned: сначала RU, потом EN
    idx = 0

    # RU
    for i, pwd in enumerate(ru_passwords, start=1):
        idx += 1
        pptx_out = work / f"ru_{i:04d}.pptx"
        render_single_brochure_pptx(template_ru, pwd, qr_png_paths[idx-1], str(pptx_out))
        pdfs.append(convert_pptx_to_pdf(soffice_bin, str(pptx_out), str(pdf_dir)))

    # EN
    for i, pwd in enumerate(en_passwords, start=1):
        idx += 1
        pptx_out = work / f"en_{i:04d}.pptx"
        render_single_brochure_pptx(template_en, pwd, qr_png_paths[idx-1], str(pptx_out))
        pdfs.append(convert_pptx_to_pdf(soffice_bin, str(pptx_out), str(pdf_dir)))

    merger = PdfMerger()
    for p in pdfs:
        merger.append(p)
    merger.write(out_pdf_path)
    merger.close()
