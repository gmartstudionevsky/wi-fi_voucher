from __future__ import annotations

import json
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
SOURCE_DIR = ROOT / "api" / "templates"
OUTPUT_DIR = ROOT / "web" / "public" / "templates"
SOFFICE = Path(r"C:\Program Files\LibreOffice\program\soffice.exe")
TOKENS = {"{{PASSWORD}}", "{{QR_WIFI}}"}


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def build_template(language: str, source: Path, work: Path) -> dict:
    presentation = Presentation(source)
    metadata: dict[str, object] = {
        "language": language,
        "slide_width": presentation.slide_width,
        "slide_height": presentation.slide_height,
        "password": None,
        "qr": None,
    }

    for slide_index, slide in enumerate(presentation.slides):
        for shape in list(iter_shapes(slide.shapes)):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if text not in TOKENS:
                continue
            key = "password" if text == "{{PASSWORD}}" else "qr"
            metadata[key] = {
                "page": slide_index,
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
            }
            remove_shape(shape)

    cleaned = work / f"brochure_{language}.pptx"
    presentation.save(cleaned)
    subprocess.run(
        [
            str(SOFFICE),
            f"-env:UserInstallation={(work / f'profile-{language}').resolve().as_uri()}",
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--norestore",
            "--convert-to",
            "pdf",
            "--outdir",
            str(work),
            str(cleaned),
        ],
        check=True,
        timeout=120,
        capture_output=True,
        text=True,
    )
    built_pdf = cleaned.with_suffix(".pdf")
    target_pdf = OUTPUT_DIR / built_pdf.name
    shutil.copy2(built_pdf, target_pdf)
    metadata["file"] = f"/templates/{target_pdf.name}"
    return metadata


def main() -> None:
    if not SOFFICE.exists():
        raise FileNotFoundError(f"LibreOffice not found: {SOFFICE}")
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="wifi-voucher-web-templates-") as temp:
        work = Path(temp)
        templates = {
            language: build_template(language, SOURCE_DIR / f"brochure_{language}.pptx", work)
            for language in ("ru", "en")
        }
    (OUTPUT_DIR / "layout.json").write_text(
        json.dumps({"templates": templates}, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
